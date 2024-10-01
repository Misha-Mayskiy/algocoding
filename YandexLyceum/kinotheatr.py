# Импорты
from datetime import datetime, timedelta

import openpyxl
from docx import Document
from openpyxl.chart import BarChart, Reference
from pptx import Presentation
from pptx.util import Inches


# Основной код
class Movie:
    def __init__(self, title, duration):
        self.title = title
        self.duration = duration


class Ticket:
    def __init__(self, session, seat):
        self.session = session
        self.seat = seat


class Session:
    def __init__(self, movie, hall, start_time):
        self.movie = movie
        self.hall = hall
        self.start_time = start_time
        self.tickets_sold = []

    def sell_ticket(self, seat):
        if seat in self.hall.seats and seat not in self.tickets_sold:
            ticket = Ticket(self, seat)
            self.tickets_sold.append(ticket)
            return ticket
        return None

    def is_seat_available(self, seat):
        return seat in self.hall.seats and seat not in self.tickets_sold


class Hall:
    def __init__(self, name, rows, seats_per_row):
        self.name = name
        self.seats = [(row, seat) for row in range(1, rows + 1) for seat in range(1, seats_per_row + 1)]
        self.sessions = []

    def add_session(self, session):
        self.sessions.append(session)


class Cinema:
    def __init__(self, name):
        self.name = name
        self.halls = []

    def add_hall(self, hall):
        self.halls.append(hall)

    def find_nearest_session(self, movie, current_time):
        available_sessions = []
        for hall in self.halls:
            for session in hall.sessions:
                if session.movie == movie and session.start_time > current_time and len(session.tickets_sold) < len(
                        hall.seats):
                    available_sessions.append(session)
        return available_sessions

    def get_hall_plan(self, session):
        plan = []
        for seat in session.hall.seats:
            if seat in session.tickets_sold:
                plan.append("X")  # Занято
            else:
                plan.append("O")  # Свободно
        return plan


# Функции для генерации документов

def generate_session_schedule(cinema, past_month):
    doc = Document()
    doc.add_heading(f'Расписание сеансов за {past_month.strftime("%B %Y")}', 0)

    for hall in cinema.halls:
        doc.add_heading(hall.name, level=1)
        for session in hall.sessions:
            session_date = datetime.strptime(session.start_time, "%Y-%m-%d %H:%M")
            if past_month <= session_date < past_month + timedelta(days=30):
                doc.add_paragraph(f'{session.movie.title} - {session.start_time}')

    doc.save(f'{cinema.name}_schedule_{past_month.strftime("%B_%Y")}.docx')


def generate_cinema_load_report(cinema, report_date):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Cinema Load"

    ws.append(["Время суток", "Количество проданных билетов"])

    time_slots = {"Утро": 0, "День": 0, "Вечер": 0, "Ночь": 0}

    for hall in cinema.halls:
        for session in hall.sessions:
            session_date = datetime.strptime(session.start_time, "%Y-%m-%d %H:%M")
            if session_date.date() == report_date:
                tickets_sold = len(session.tickets_sold)
                hour = session_date.hour
                if 6 <= hour < 12:
                    time_slots["Утро"] += tickets_sold
                elif 12 <= hour < 18:
                    time_slots["День"] += tickets_sold
                elif 18 <= hour < 24:
                    time_slots["Вечер"] += tickets_sold
                else:
                    time_slots["Ночь"] += tickets_sold

    for time, count in time_slots.items():
        ws.append([time, count])

    # Создаем график
    chart = BarChart()
    data = Reference(ws, min_col=2, min_row=1, max_row=5)
    categories = Reference(ws, min_col=1, min_row=2, max_row=5)
    chart.add_data(data, titles_from_data=True)
    chart.set_categories(categories)
    chart.title = "Загруженность кинотеатра"
    ws.add_chart(chart, "E5")

    wb.save(f"{cinema.name}_load_report_{report_date}.xlsx")


def generate_movie_brochure(cinema, new_movies):
    prs = Presentation()

    for movie in new_movies:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = f"Рекламный буклет: {movie.title}"

        subtitle = slide.shapes.placeholders[1]
        subtitle.text = f"Продолжительность: {movie.duration} мин"

        slide_content = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(2))
        content = slide_content.text_frame
        content.text = "Сеансы в кинотеатрах:"

        for hall in cinema.halls:
            for session in hall.sessions:
                if session.movie == movie:
                    p = content.add_paragraph()
                    p.text = f"{hall.name}: {session.start_time}"

    prs.save(f"Movie_Brochure_{datetime.now().strftime('%Y_%m_%d')}.pptx")


# Пример использования
cinema = Cinema("Cinemax")
hall1 = Hall("Hall 1", 5, 10)
cinema.add_hall(hall1)
movie1 = Movie("Inception", 148)
session1 = Session(movie1, hall1, "2024-10-01 19:00")
hall1.add_session(session1)
ticket = session1.sell_ticket((1, 1))

# Тестирование функций
past_month = datetime(2024, 9, 1)
generate_session_schedule(cinema, past_month)

report_date = datetime(2024, 10, 1).date()
generate_cinema_load_report(cinema, report_date)

new_movies = [movie1]
generate_movie_brochure(cinema, new_movies)
